from flask import Flask, render_template, request, send_file, jsonify
import sqlite3
import pandas as pd
from datetime import datetime, timedelta
import io
import json
import subprocess
import os
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB max
DB_PATH = "bump_data.db"
CONFIG_STATIONS_PATH = "config_stations.xlsx"

def load_station_config():
    """Charge la configuration des stations depuis le fichier Excel"""
    try:
        # Lire la première feuille quelle que soit son nom
        df = pd.read_excel(CONFIG_STATIONS_PATH, sheet_name=0)
        
        # Mapping des stations par nom (car pas de station_id dans l'Excel)
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        cursor.execute("SELECT id, name FROM stations")
        stations_db = {row[1]: str(row[0]) for row in cursor.fetchall()}
        conn.close()
        
        # Convertir en dictionnaire avec station_id comme clé
        config = {}
        for _, row in df.iterrows():
            station_name = str(row['station_name'])
            # Trouver l'ID correspondant dans la base
            station_id = None
            for db_name, db_id in stations_db.items():
                if station_name in db_name or db_name in station_name:
                    station_id = db_id
                    break
            
            if station_id:
                date_ouv = str(row['date_ouverture'])
                # Extraire juste la date (AAAA-MM-JJ)
                if ' ' in date_ouv:
                    date_ouv = date_ouv.split(' ')[0]
                
                config[station_id] = {
                    'nb_places': int(row['nb_places']),
                    'bp_td': float(row['nb_recharge_pdc_bp_td']),  # Nom correct de la colonne
                    'date_ouverture': date_ouv[:10],
                    'commentaire': ''
                }
        
        print(f"✅ Configuration chargée pour {len(config)} stations")
        return config
    except FileNotFoundError:
        print(f"âš ï¸  Fichier {CONFIG_STATIONS_PATH} non trouvé. Utilisation des valeurs par défaut.")
        return {}
    except Exception as e:
        print(f"âš ï¸  Erreur lors du chargement de {CONFIG_STATIONS_PATH}: {e}")
        import traceback
        traceback.print_exc()
        return {}

def get_station_config(station_id):
    """Récupère la config d'une station, ou valeurs par défaut"""
    config = load_station_config()
    return config.get(str(station_id), {
        'nb_places': 2,
        'bp_td': 2.0,
        'date_ouverture': '2024-01-01',
        'commentaire': ''
    })

def get_stations():
    """Récupère la liste des stations AVIA VOLT (Bump + ETL)"""
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()

    # Stations Bump
    cursor.execute("""
        SELECT DISTINCT s.id, s.name 
        FROM stations s
        JOIN chargers c ON s.id = c.location_id
        JOIN charging_sessions cs ON c.id = cs.charger_id
    """)
    bump_stations = cursor.fetchall()

    # Stations ETL (sessions réussies avec énergie > 0)
    cursor.execute("""
        SELECT DISTINCT 'ETL-' || station_name as id, station_name as name
        FROM etl_sessions
        WHERE energie_wh > 0 AND charge_reussie = 1
    """)
    etl_stations = cursor.fetchall()

    conn.close()

    all_stations = list(bump_stations) + list(etl_stations)
    all_stations.sort(key=lambda x: x[1])
    return all_stations


def is_etl_station(station_id):
    """Détermine si une station est de type ETL"""
    return str(station_id).startswith('ETL-')

def get_etl_station_name(station_id):
    """Extrait le nom de station depuis l'ID ETL"""
    return str(station_id)[4:]

def get_station_display_name(station_id):
    """Retourne le nom d'affichage d'une station (Bump ou ETL)"""
    if is_etl_station(station_id):
        return get_etl_station_name(station_id)
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM stations WHERE id = ?", [station_id])
    row = cursor.fetchone()
    conn.close()
    return row[0] if row else f"Station {station_id}"

def generate_graph_data(station_id, date_debut, date_fin):
    """Génère les données pour le graphique d'une station (Bump ou ETL)"""
    conn = sqlite3.connect(DB_PATH)

    station_config = get_station_config(station_id)
    BP_TD_VALUE = station_config['bp_td']
    NB_PLACES = station_config['nb_places']
    DATE_OUVERTURE = station_config['date_ouverture']

    date_debut_str = str(date_debut)
    date_fin_str = str(date_fin)
    date_ouverture_str = str(DATE_OUVERTURE)[:10]

    if date_debut_str < date_ouverture_str:
        date_debut_ajustee = date_ouverture_str
    else:
        date_debut_ajustee = date_debut_str

    # === REQUÊTE selon le type ===
    if is_etl_station(station_id):
        etl_name = get_etl_station_name(station_id)
        query = """
            SELECT 
                DATE(date_debut) as date,
                energie_wh / 1000.0 as energy_kwh
            FROM etl_sessions
            WHERE station_name = ?
            AND charge_reussie = 1
            AND energie_wh > 0
            AND DATE(date_debut) BETWEEN ? AND ?
            ORDER BY date_debut
        """
        df = pd.read_sql_query(query, conn, params=[etl_name, date_debut_ajustee, date_fin_str])
    else:
        query = """
            SELECT 
                DATE(cs.start_time) as date,
                cs.energy_kwh
            FROM charging_sessions cs
            JOIN chargers c ON cs.charger_id = c.id
            WHERE c.location_id = ?
            AND DATE(cs.start_time) BETWEEN ? AND ?
            AND cs.energy_kwh > 0
            ORDER BY cs.start_time
        """
        df = pd.read_sql_query(query, conn, params=[station_id, date_debut_ajustee, date_fin_str])

    conn.close()

    if not df.empty:
        df['date'] = pd.to_datetime(df['date'])

    # Index complet
    date_range = pd.date_range(start=date_debut_ajustee, end=date_fin, freq='D')
    all_dates = pd.DataFrame(index=date_range)
    all_dates.index.name = 'date'

    if df.empty:
        daily_stats = all_dates.copy()
        daily_stats['nb_recharges'] = 0
        daily_stats['kwh_total'] = 0.0
        daily_stats['nb_recharges_par_pdc'] = 0.0
        daily_stats['kwh_par_pdc'] = 0.0
        daily_stats['moyenne_glissante'] = 0.0
        daily_stats['bp_td'] = BP_TD_VALUE
    else:
        daily_stats = df.groupby('date').agg({
            'energy_kwh': ['count', 'sum']
        })
        daily_stats.columns = ['nb_recharges', 'kwh_total']
        daily_stats = all_dates.join(daily_stats, how='left')

        daily_stats['nb_recharges'] = daily_stats['nb_recharges'].fillna(0).astype(int)
        daily_stats['kwh_total'] = daily_stats['kwh_total'].fillna(0)
        daily_stats['nb_recharges_par_pdc'] = (daily_stats['nb_recharges'] / NB_PLACES).round(2)
        daily_stats['kwh_par_pdc'] = (daily_stats['kwh_total'] / NB_PLACES).round(2)
        daily_stats['moyenne_glissante'] = daily_stats['nb_recharges_par_pdc'].rolling(window=14, min_periods=1).mean().round(2)
        daily_stats['bp_td'] = BP_TD_VALUE

    daily_stats['is_weekend'] = daily_stats.index.dayofweek.isin([5, 6])

    # Préparer les données pour Chart.js
    dates = [d.strftime('%Y-%m-%d') for d in daily_stats.index.tolist()]
    nb_recharges_par_pdc = [float(x) for x in daily_stats['nb_recharges_par_pdc'].tolist()]
    nb_recharges_total = [float(x) for x in daily_stats['nb_recharges'].tolist()]
    kwh_par_pdc = [float(x) for x in daily_stats['kwh_par_pdc'].tolist()]
    kwh_total = [float(x) for x in daily_stats['kwh_total'].tolist()]

    # Cumuls
    kwh_par_pdc_cumul = []
    kwh_total_cumul = []
    cumul_pdc = 0
    cumul_total = 0
    for i in range(len(kwh_par_pdc)):
        cumul_pdc += kwh_par_pdc[i]
        cumul_total += kwh_total[i]
        kwh_par_pdc_cumul.append(round(cumul_pdc, 2))
        kwh_total_cumul.append(round(cumul_total, 2))

    moyenne_glissante = [float(x) for x in daily_stats['moyenne_glissante'].tolist()]
    moyenne_glissante_total = daily_stats['nb_recharges'].rolling(window=14, min_periods=1).mean().round(2).tolist()
    moyenne_glissante_total = [float(x) for x in moyenne_glissante_total]
    bp_td = [float(x) for x in daily_stats['bp_td'].tolist()]
    bp_td_total = [float(BP_TD_VALUE * NB_PLACES) for _ in range(len(dates))]
    is_weekend = [bool(x) for x in daily_stats['is_weekend'].tolist()]

    return {
        'dates': dates,
        'nb_recharges_par_pdc': nb_recharges_par_pdc,
        'nb_recharges_total': nb_recharges_total,
        'kwh_par_pdc': kwh_par_pdc,
        'kwh_total': kwh_total,
        'kwh_par_pdc_cumul': kwh_par_pdc_cumul,
        'kwh_total_cumul': kwh_total_cumul,
        'moyenne_glissante': moyenne_glissante,
        'moyenne_glissante_total': moyenne_glissante_total,
        'bp_td': bp_td,
        'bp_td_total': bp_td_total,
        'is_weekend': is_weekend,
        'nb_places': NB_PLACES
    }

def calculate_stats(station_id, date_debut, date_fin):
    """Calcule les statistiques pour une station sur une période"""
    conn = sqlite3.connect(DB_PATH)
    
    # Récupérer la configuration de la station
    station_config = get_station_config(station_id)
    NB_PLACES = station_config['nb_places']
    DATE_OUVERTURE = station_config['date_ouverture']
    
    # Convertir les dates en string format AAAA-MM-JJ pour comparaison
    date_debut_str = str(date_debut)
    date_fin_str = str(date_fin)
    date_ouverture_str = str(DATE_OUVERTURE)[:10]  # Garder seulement AAAA-MM-JJ
    
    # Ajuster la date de début si elle est antérieure Ã  l'ouverture (comparaison string)
    if date_debut_str < date_ouverture_str:
        date_debut_ajustee = date_ouverture_str
        print(f"⚠️  Station {station_id}: date demandée ({date_debut_str}) avant ouverture ({date_ouverture_str}), ajusté à {date_debut_ajustee}")
    else:
        date_debut_ajustee = date_debut_str
        print(f"✓ Station {station_id}: période {date_debut_str} à {date_fin_str} OK")
    
    # Récupérer les sessions de la période (ajustée)
    # === REQUÊTE selon le type de station ===
    if is_etl_station(station_id):
        etl_name = get_etl_station_name(station_id)
        query = """
            SELECT 
                DATE(date_debut) as date,
                id_transaction as id,
                date_debut as start_time,
                date_fin as end_time,
                energie_wh / 1000.0 as energy_kwh,
                'ETL Charger' as charger_name
            FROM etl_sessions
            WHERE station_name = ?
            AND charge_reussie = 1
            AND energie_wh > 0
            AND DATE(date_debut) BETWEEN ? AND ?
            ORDER BY date_debut
        """
        df = pd.read_sql_query(query, conn, params=[etl_name, date_debut_ajustee, date_fin_str])
    else:
        query = """
            SELECT 
                DATE(cs.start_time) as date,
                cs.id,
                cs.start_time,
                cs.end_time,
                cs.energy_kwh,
                c.name as charger_name
            FROM charging_sessions cs
            JOIN chargers c ON cs.charger_id = c.id
            WHERE c.location_id = ?
            AND DATE(cs.start_time) BETWEEN ? AND ?
            AND cs.energy_kwh > 0
            ORDER BY cs.start_time
        """
        df = pd.read_sql_query(query, conn, params=[station_id, date_debut_ajustee, date_fin_str])
    
    conn.close()
    
    print(f"   📊 DataFrame: {len(df)} lignes")
    
    if df.empty:
        return None
    
    # Convertir les dates
    df['start_time'] = pd.to_datetime(df['start_time'])
    df['end_time'] = pd.to_datetime(df['end_time'])
    df['date'] = pd.to_datetime(df['date']).dt.date  # Convertir explicitement en date
    
    # Calculer la durée en heures
    df['duree_h'] = (df['end_time'] - df['start_time']).dt.total_seconds() / 3600
    
    # Identifier le type de prise (simplifié - Ã  adapter selon vos données)
    # Pour l'instant, on suppose CCS si >50kW, sinon T2
    df['type_prise'] = df['energy_kwh'].apply(lambda x: 'CCS' if x > 50 else 'T2')
    
    # Créer un index complet avec toutes les dates de la période (depuis l'ouverture)
    date_range = pd.date_range(start=date_debut_ajustee, end=date_fin, freq='D')
    all_dates_index = date_range.date
    all_dates = pd.DataFrame(index=all_dates_index)
    all_dates.index.name = 'date'
    
    # Stats par jour
    stats_jour = df.groupby('date').agg({
        'id': 'count',  # Nombre de recharges
        'energy_kwh': ['sum', 'mean'],  # Total et moyenne kWh
        'duree_h': 'mean'  # Durée moyenne
    }).round(2)
    
    stats_jour.columns = ['Nombre de recharges', 'Total kWh', 'kWh moyen', 'Durée moyenne (h)']
    
    # Fusionner avec toutes les dates pour inclure les jours Ã  0
    stats_jour = all_dates.join(stats_jour, how='left')
    
    # Remplir les valeurs manquantes par 0
    stats_jour['Nombre de recharges'] = stats_jour['Nombre de recharges'].fillna(0).astype(int)
    stats_jour['Total kWh'] = stats_jour['Total kWh'].fillna(0).round(2)
    stats_jour['kWh moyen'] = stats_jour['kWh moyen'].fillna(0).round(2)
    stats_jour['Durée moyenne (h)'] = stats_jour['Durée moyenne (h)'].fillna(0).round(2)
    
    # Trier par date CROISSANTE (plus ancien en haut)
    stats_jour = stats_jour.sort_index(ascending=True)
    
    # CA prévisionnel (0.59€/kWh par exemple - Ã  ajuster)
    PRIX_KWH_TTC = 0.59
    stats_jour['CA HT (€)'] = (stats_jour['Total kWh'] * PRIX_KWH_TTC / 1.2).round(2)
    
    # Recharges par place par jour (NB_PLACES vient de la config)
    stats_jour['Recharges/place/jour'] = (stats_jour['Nombre de recharges'] / NB_PLACES).round(2)
    
    # Moyenne glissante 15 jours
    stats_jour['Moy. glissante 15j'] = stats_jour['Recharges/place/jour'].rolling(window=15, min_periods=1).mean().round(2)
    
    # Durée d'occupation moyenne par place (total durée / nb places)
    duree_occupation_jour = df.groupby('date')['duree_h'].sum() / NB_PLACES
    stats_jour['Occupation/place (h)'] = duree_occupation_jour.round(2)
    
    # Compter CCS et T2 par jour (mais ne pas les inclure dans le tableau final)
    # type_prise_counts = df.groupby(['date', 'type_prise']).size().unstack(fill_value=0)
    # stats_jour['CCS'] = type_prise_counts.get('CCS', 0)
    # stats_jour['T2'] = type_prise_counts.get('T2', 0)
    
    # Ligne total
    # Calculer le nombre de recharges total (uniquement les jours avec recharge)
    total_recharges = stats_jour['Nombre de recharges'].sum()
    # Pour les moyennes, calculer seulement sur les jours avec activité
    jours_avec_activite = stats_jour[stats_jour['Nombre de recharges'] > 0]
    
    total = pd.DataFrame({
        'Nombre de recharges': [total_recharges],
        'Total kWh': [stats_jour['Total kWh'].sum()],
        'kWh moyen': [jours_avec_activite['kWh moyen'].mean() if len(jours_avec_activite) > 0 else 0],
        'Durée moyenne (h)': [jours_avec_activite['Durée moyenne (h)'].mean() if len(jours_avec_activite) > 0 else 0],
        'CA HT (€)': [stats_jour['CA HT (€)'].sum()],
        'Recharges/place/jour': [stats_jour['Recharges/place/jour'].mean()],
        'Moy. glissante 15j': [stats_jour['Moy. glissante 15j'].mean()],
        'Occupation/place (h)': [stats_jour['Occupation/place (h)'].mean()]
    }, index=['TOTAL']).round(2)
    
    stats_jour = pd.concat([stats_jour, total])
    
    # Renommer l'index pour avoir "Date" comme nom de colonne
    stats_jour.index.name = 'Date'
    
    return stats_jour, NB_PLACES

def calculate_recap(station_ids, date_debut, date_fin):
    """Calcule le tableau récapitulatif pour plusieurs stations"""
    conn = sqlite3.connect(DB_PATH)
    
    recap_data = []
    
    for station_id in station_ids:
        # Infos station - utiliser la fonction qui gère Bump et ETL
        station_name = get_station_display_name(station_id)
        print(f"   🔍 Recap: {station_name}")
        
        # Récupérer la configuration de la station
        station_config = get_station_config(station_id)
        NB_PLACES = station_config['nb_places']
        DATE_OUVERTURE = station_config['date_ouverture']
        
        # Convertir les dates en string format AAAA-MM-JJ pour comparaison
        date_debut_str = str(date_debut)
        date_fin_str = str(date_fin)
        date_ouverture_str = str(DATE_OUVERTURE)[:10]
        
        # Ajuster la date de début si elle est antérieure Ã  l'ouverture
        if date_debut_str < date_ouverture_str:
            date_debut_ajustee = date_ouverture_str
        else:
            date_debut_ajustee = date_debut_str
        
        # Stats - selon le type de station
        if is_etl_station(station_id):
            etl_name = get_etl_station_name(station_id)
            query = """
                SELECT 
                    id_transaction as id,
                    date_debut as start_time,
                    date_fin as end_time,
                    energie_wh / 1000.0 as energy_kwh
                FROM etl_sessions
                WHERE station_name = ?
                AND charge_reussie = 1
                AND energie_wh > 0
                AND DATE(date_debut) BETWEEN ? AND ?
            """
            df = pd.read_sql_query(query, conn, params=[etl_name, date_debut_ajustee, date_fin_str])
        else:
            query = """
                SELECT 
                    cs.id,
                    cs.start_time,
                    cs.end_time,
                    cs.energy_kwh
                FROM charging_sessions cs
                JOIN chargers c ON cs.charger_id = c.id
                WHERE c.location_id = ?
                AND DATE(cs.start_time) BETWEEN ? AND ?
                AND cs.energy_kwh > 0
            """
            df = pd.read_sql_query(query, conn, params=[station_id, date_debut_ajustee, date_fin_str])
        
        if df.empty:
            continue
        
        df['start_time'] = pd.to_datetime(df['start_time'])
        df['end_time'] = pd.to_datetime(df['end_time'])
        df['duree_h'] = (df['end_time'] - df['start_time']).dt.total_seconds() / 3600
        df['type_prise'] = df['energy_kwh'].apply(lambda x: 'CCS' if x > 50 else 'T2')
        
        PRIX_KWH_TTC = 0.59
        nb_jours = (pd.to_datetime(date_fin) - pd.to_datetime(date_debut_ajustee)).days + 1
        
        recap_data.append({
            'Station': station_name,
            'Nombre de places': NB_PLACES,
            'BP TD': station_config['bp_td'],
            'Période': f"{date_debut} au {date_fin}",
            'Nombre de recharges': len(df),
            'Total kWh': round(df['energy_kwh'].sum(), 2),
            'kWh moyen par recharge': round(df['energy_kwh'].mean(), 2),
            'CA HT (€)': round(df['energy_kwh'].sum() * PRIX_KWH_TTC / 1.2, 2),
            'Nb de recharge par place par jour': round(len(df) / NB_PLACES / nb_jours, 2),
            'Durée moyenne d\'une recharge (h)': round(df['duree_h'].mean(), 2),
            'Durée d\'occupation moyenne par place (h)': round(df['duree_h'].sum() / NB_PLACES / nb_jours, 2)
        })
    
    conn.close()
    
    return pd.DataFrame(recap_data)

@app.route('/')
def index():
    stations = get_stations()
    return render_template('index.html', stations=stations)

@app.route('/graphiques', methods=['POST'])
def graphiques():
    date_debut = request.form.get('date_debut')
    date_fin = request.form.get('date_fin')
    station_ids = request.form.getlist('stations')
    
    if not station_ids:
        return "Veuillez sélectionner au moins une station", 400
    
    # Générer les graphiques pour chaque station
    graphiques_stations = {}
    for station_id in station_ids:
        graphique_data = generate_graph_data(station_id, date_debut, date_fin)
        if graphique_data:
            graphiques_stations[station_id] = {
                'name': get_station_display_name(station_id),
                'data': graphique_data
            }
    
    return render_template('graphiques.html',
                         graphiques_stations=graphiques_stations,
                         date_debut=date_debut,
                         date_fin=date_fin,
                         station_ids=station_ids)

@app.route('/generer', methods=['POST'])
def generer():
    date_debut = request.form.get('date_debut')
    date_fin = request.form.get('date_fin')
    station_ids = request.form.getlist('stations')
    
    if not station_ids:
        return "Veuillez sélectionner au moins une station", 400
    
    # Calculer les stats pour chaque station
    stats_stations = {}
    station_ids_with_data = []
    for station_id in station_ids:
        print(f"🔍 Traitement de {station_id}...")
        result = calculate_stats(station_id, date_debut, date_fin)
        print(f"   Résultat: {result is not None}")
        if result is not None:
            stats, nb_places = result
            print(f"   ✅ Stats OK, {len(stats)} lignes, {nb_places} places")
            stats_stations[station_id] = {
                'name': get_station_display_name(station_id),
                'stats': stats.to_html(classes='table table-bordered', border=1),
                'nb_places': nb_places
            }
            station_ids_with_data.append(station_id)
        else:
            print(f"   ❌ Aucune donnée retournée")
    
    # Tableau récapitulatif
    recap = calculate_recap(station_ids, date_debut, date_fin)
    recap_html = recap.to_html(classes='table table-bordered', index=False, border=1) if not recap.empty else None
    
    return render_template('resultats.html', 
                         stats_stations=stats_stations,
                         recap_html=recap_html,
                         date_debut=date_debut,
                         date_fin=date_fin,
                         station_ids_with_data=station_ids_with_data)

@app.route('/exporter', methods=['POST'])
def exporter():
    date_debut = request.form.get('date_debut')
    date_fin = request.form.get('date_fin')
    station_ids = request.form.getlist('station_ids')  # getlist au lieu de get
    
    # Créer un fichier Excel avec plusieurs feuilles
    output = io.BytesIO()
    writer = pd.ExcelWriter(output, engine='xlsxwriter')
    
    # Feuille pour chaque station
    for station_id in station_ids:
        result = calculate_stats(station_id, date_debut, date_fin)
        if result is not None:
            stats, nb_places = result
            sheet_name = get_station_display_name(station_id)[:31]  # Limite Excel
            stats.to_excel(writer, sheet_name=sheet_name)
    
    # Feuille récapitulative
    recap = calculate_recap(station_ids, date_debut, date_fin)
    if not recap.empty:
        recap.to_excel(writer, sheet_name='Récapitulatif', index=False)
    
    writer.close()
    output.seek(0)
    
    return send_file(output, 
                    mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                    as_attachment=True,
                    download_name=f'export_stations_{date_debut}_{date_fin}.xlsx')

@app.route('/generer_rapport_word', methods=['POST'])
def generer_rapport_word():
    """Génère un rapport Word avec tableaux et graphiques capturés depuis le navigateur"""
    import base64
    
    date_debut = request.form.get('date_debut')
    date_fin = request.form.get('date_fin')
    station_ids = request.form.getlist('station_ids')
    
    # Récupérer les images capturées depuis le navigateur
    graphs_data_json = request.form.get('graphs_data', '{}')
    graphs_data = json.loads(graphs_data_json)
    
    if not station_ids:
        return "Veuillez sélectionner au moins une station", 400
    
    # Créer le document Word
    doc = Document()
    
    # Titre principal
    title = doc.add_heading('Rapport AVIA VOLT - Statistiques des Stations', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Période
    period = doc.add_paragraph()
    period.add_run(f'Période : {date_debut} au {date_fin}').bold = True
    period.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    doc.add_paragraph()  # Espace
    
    # Récupérer les noms des stations
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    placeholders = ','.join('?' * len(station_ids))
    cursor.execute(f"SELECT id, name FROM stations WHERE id IN ({placeholders})", station_ids)
    stations = {str(row[0]): row[1] for row in cursor.fetchall()}
    conn.close()
    
    # Pour chaque station
    for station_id in station_ids:
        station_name = stations.get(str(station_id), f"Station {station_id}")
        
        # Titre de la station
        doc.add_heading(f'Station: {station_name}', 1)
        doc.add_paragraph()
        
        # Calculer les stats
        result = calculate_stats(station_id, date_debut, date_fin)
        
        if result is not None:
            stats_jour, NB_PLACES = result  # Décompresser le tuple
            # Configuration de la station
            station_config = get_station_config(station_id)
            NB_PLACES = station_config['nb_places']
            BP_TD = station_config['bp_td']
            DATE_OUVERTURE = station_config['date_ouverture']
            
            doc.add_heading('Configuration', 2)
            config_text = f"""Nombre de points de charge: {NB_PLACES}
Business Plan - Taux de remplissage journalier: {BP_TD} recharges/PDC
Date d'ouverture: {DATE_OUVERTURE}"""
            doc.add_paragraph(config_text)
            doc.add_paragraph()
            
            # Tableau récapitulatif (uniquement la ligne TOTAL)
            doc.add_heading('Tableau Récapitulatif', 2)
            
            # Sélectionner uniquement la ligne TOTAL
            total_row = stats_jour.iloc[-1]
            
            table = doc.add_table(rows=2, cols=6)
            table.style = 'Light Grid Accent 1'
            
            # Headers
            headers = ['Recharges', 'Total kWh', 'CA HT (€)', 'Recharges/PDC/j', 'Moy. 15j', 'Occup/PDC (h)']
            for i, header in enumerate(headers):
                cell = table.rows[0].cells[i]
                cell.text = header
                cell.paragraphs[0].runs[0].font.bold = True
            
            # Valeurs TOTAL
            values = [
                str(int(total_row['Nombre de recharges'])),
                str(total_row['Total kWh']),
                str(total_row['CA HT (€)']),
                str(total_row['Recharges/place/jour']),
                str(total_row['Moy. glissante 15j']),
                str(total_row['Occupation/place (h)'])
            ]
            for i, value in enumerate(values):
                table.rows[1].cells[i].text = value
            
            doc.add_paragraph()
            
            # Ajouter les graphiques capturés depuis le navigateur
            doc.add_heading('Graphiques d\'Analyse', 2)
            
            # 1. Graphique PDC - Vue Jour
            doc.add_heading('Vue Jour - Par Point de Charge', 3)
            img_key = f'{station_id}_pdc_jour'
            if img_key in graphs_data:
                img_data = graphs_data[img_key].split(',')[1]
                img_bytes = base64.b64decode(img_data)
                img_buffer = io.BytesIO(img_bytes)
                doc.add_picture(img_buffer, width=Inches(6.5))
            else:
                doc.add_paragraph("Graphique non disponible")
            
            # 2. Graphique Total - Vue Jour
            doc.add_heading('Vue Jour - Global Station', 3)
            img_key = f'{station_id}_total_jour'
            if img_key in graphs_data:
                img_data = graphs_data[img_key].split(',')[1]
                img_bytes = base64.b64decode(img_data)
                img_buffer = io.BytesIO(img_bytes)
                doc.add_picture(img_buffer, width=Inches(6.5))
            else:
                doc.add_paragraph("Graphique non disponible")
            
            # 3. Graphique Total - Vue Mois
            doc.add_heading('Vue Mois - Global Station', 3)
            img_key = f'{station_id}_total_mois'
            if img_key in graphs_data:
                img_data = graphs_data[img_key].split(',')[1]
                img_bytes = base64.b64decode(img_data)
                img_buffer = io.BytesIO(img_bytes)
                doc.add_picture(img_buffer, width=Inches(6.5))
            else:
                doc.add_paragraph("Graphique non disponible")
        
        else:
            doc.add_paragraph("Aucune donnée disponible pour cette station sur la période sélectionnée.")
    
    # Sauvegarder le document
    output = io.BytesIO()
    doc.save(output)
    output.seek(0)
    
    return send_file(output,
                    mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    as_attachment=True,
                    download_name=f'rapport_avia_volt_{date_debut}_{date_fin}.docx')

@app.route('/generer_rapport_pptx', methods=['POST'])
def generer_rapport_pptx():
    """Génère un rapport PowerPoint avec tableaux et graphiques capturés depuis le navigateur"""
    import base64
    
    date_debut = request.form.get('date_debut')
    date_fin = request.form.get('date_fin')
    station_ids = request.form.getlist('station_ids')
    
    # Récupérer les images capturées depuis le navigateur
    graphs_data_json = request.form.get('graphs_data', '{}')
    graphs_data = json.loads(graphs_data_json)
    
    if not station_ids:
        return "Veuillez sélectionner au moins une station", 400
    
    # Créer la présentation PowerPoint
    prs = Presentation()
    prs.slide_width = PptxInches(10)  # Format paysage
    prs.slide_height = PptxInches(7.5)
    
    # Slide de titre
    slide_layout = prs.slide_layouts[0]  # Layout titre
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Rapport AVIA VOLT - Statistiques des Stations"
    subtitle.text = f"Période : {date_debut} au {date_fin}"
    
    # Récupérer les noms des stations
    conn = sqlite3.connect(DB_PATH)
    cursor = conn.cursor()
    placeholders = ','.join('?' * len(station_ids))
    cursor.execute(f"SELECT id, name FROM stations WHERE id IN ({placeholders})", station_ids)
    stations = {str(row[0]): row[1] for row in cursor.fetchall()}
    conn.close()
    
    # Pour chaque station
    for station_id in station_ids:
        station_name = stations.get(str(station_id), f"Station {station_id}")
        
        # Calculer les stats
        result = calculate_stats(station_id, date_debut, date_fin)
        
        if result is not None:
            stats_jour, NB_PLACES = result
            
            # Configuration de la station
            station_config = get_station_config(station_id)
            BP_TD = station_config['bp_td']
            DATE_OUVERTURE = station_config['date_ouverture']
            
            # === SLIDE 1 : Tableau récapitulatif détaillé ===
            slide_layout = prs.slide_layouts[5]  # Layout vide
            slide = prs.slides.add_slide(slide_layout)
            
            # Titre
            title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"{station_name} - Tableau Récapitulatif"
            title_para.font.size = Pt(28)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Tableau récapitulatif détaillé
            total_row = stats_jour.iloc[-1]
            
            # Position et taille du tableau (11 colonnes)
            left = PptxInches(0.3)
            top = PptxInches(1.5)
            width = PptxInches(9.4)
            height = PptxInches(2)
            
            # Créer tableau 2 lignes x 11 colonnes
            table = slide.shapes.add_table(2, 11, left, top, width, height).table
            
            # Headers
            headers = [
                'Station',
                'Nb places',
                'BP TD',
                'Période',
                'Nb recharges',
                'Total kWh',
                'kWh moy/rech',
                'CA HT (€)',
                'Rech/PDC/j',
                'Durée moy rech (h)',
                'Occ/PDC (h)'
            ]
            
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                # Ajuster la largeur des colonnes
                if i == 0:  # Station
                    table.columns[i].width = PptxInches(1.2)
                elif i == 3:  # Période
                    table.columns[i].width = PptxInches(1.2)
                else:
                    table.columns[i].width = PptxInches(0.75)
            
            # Valeurs
            values = [
                station_name,
                str(NB_PLACES),
                str(BP_TD),
                f"{date_debut} au {date_fin}",
                str(int(total_row['Nombre de recharges'])),
                f"{total_row['Total kWh']:.1f}",
                f"{total_row['kWh moyen']:.2f}",
                f"{total_row['CA HT (€)']:.2f}",
                f"{total_row['Recharges/place/jour']:.2f}",
                f"{total_row['Durée moyenne (h)']:.2f}",
                f"{total_row['Occupation/place (h)']:.2f}"
            ]
            
            for i, value in enumerate(values):
                cell = table.cell(1, i)
                cell.text = value
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # === SLIDE 2 : 2 Graphiques côte Ã  côte ===
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            
            # Titre
            title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2), PptxInches(9), PptxInches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"{station_name} - Analyse Graphique"
            title_para.font.size = Pt(24)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Graphique 1 : PDC Vue Jour (Ã  gauche)
            img_key = f'{station_id}_pdc_jour'
            if img_key in graphs_data:
                # Sous-titre
                subtitle_box = slide.shapes.add_textbox(PptxInches(0.3), PptxInches(0.9), PptxInches(4.5), PptxInches(0.3))
                subtitle_frame = subtitle_box.text_frame
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.text = "Vue Jour - Par Point de Charge"
                subtitle_para.font.size = Pt(14)
                subtitle_para.font.bold = True
                subtitle_para.alignment = PP_ALIGN.CENTER
                
                # Graphique
                img_data = graphs_data[img_key].split(',')[1]
                img_bytes = base64.b64decode(img_data)
                img_buffer = io.BytesIO(img_bytes)
                slide.shapes.add_picture(img_buffer, PptxInches(0.3), PptxInches(1.3), width=PptxInches(4.5))
            
            # Graphique 2 : Total Vue Mois (Ã  droite)
            img_key = f'{station_id}_total_mois'
            if img_key in graphs_data:
                # Sous-titre
                subtitle_box = slide.shapes.add_textbox(PptxInches(5.2), PptxInches(0.9), PptxInches(4.5), PptxInches(0.3))
                subtitle_frame = subtitle_box.text_frame
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.text = "Vue Mois - Global Station"
                subtitle_para.font.size = Pt(14)
                subtitle_para.font.bold = True
                subtitle_para.alignment = PP_ALIGN.CENTER
                
                # Graphique
                img_data = graphs_data[img_key].split(',')[1]
                img_bytes = base64.b64decode(img_data)
                img_buffer = io.BytesIO(img_bytes)
                slide.shapes.add_picture(img_buffer, PptxInches(5.2), PptxInches(1.3), width=PptxInches(4.5))
    
    # Sauvegarder la présentation
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    
    return send_file(output,
                    mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                    as_attachment=True,
                    download_name=f'rapport_avia_volt_{date_debut}_{date_fin}.pptx')

@app.route('/update_data', methods=['POST'])
def update_data():
    """Exécute les scripts d'import pour mettre à jour les données"""
    try:
        results = []
        
        # Liste des scripts à exécuter
        scripts = []
        
        # Vérifier si bump_import.py existe
        if os.path.exists('bump_import.py'):
            scripts.append(('bump_import.py', 'Bump'))
        
        # Vérifier si EC_import.py existe
        if os.path.exists('EC_import.py'):
            scripts.append(('EC_import.py', 'EC'))
        
        # Vérifier si etl_import.py existe
        if os.path.exists('etl_import.py'):
            scripts.append(('etl_import.py', 'ETL'))
        
        if not scripts:
            return jsonify({
                'success': False,
                'message': 'Aucun script d\'import trouvé'
            }), 404
        
        # Exécuter chaque script
        for script_name, source_name in scripts:
            try:
                result = subprocess.run(
                    ['python3', script_name],
                    capture_output=True,
                    text=True,
                    timeout=300  # 5 minutes max par script
                )
                
                if result.returncode == 0:
                    results.append({
                        'source': source_name,
                        'status': 'success',
                        'message': f'Import {source_name} réussi',
                        'output': result.stdout
                    })
                else:
                    results.append({
                        'source': source_name,
                        'status': 'error',
                        'message': f'Erreur lors de l\'import {source_name}',
                        'output': result.stderr
                    })
            except subprocess.TimeoutExpired:
                results.append({
                    'source': source_name,
                    'status': 'error',
                    'message': f'Timeout lors de l\'import {source_name} (>5min)',
                    'output': ''
                })
            except Exception as e:
                results.append({
                    'source': source_name,
                    'status': 'error',
                    'message': f'Erreur inattendue pour {source_name}',
                    'output': str(e)
                })
        
        # Vérifier si au moins un import a réussi
        success = any(r['status'] == 'success' for r in results)
        
        return jsonify({
            'success': success,
            'results': results,
            'message': 'Mise à jour terminée' if success else 'Échec de la mise à jour'
        })
        
    except Exception as e:
        return jsonify({
            'success': False,
            'message': f'Erreur lors de la mise à jour: {str(e)}'
        }), 500

if __name__ == '__main__':
    import os
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)